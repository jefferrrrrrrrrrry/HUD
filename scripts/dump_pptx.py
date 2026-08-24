import sys
from pptx import Presentation
from pptx.util import Emu

path = '周颖-开题-0830.pptx'
prs = Presentation(path)
print('slides:', len(prs.slides), '| size:', prs.slide_width, 'x', prs.slide_height,
      f'({Emu(prs.slide_width).inches:.2f} x {Emu(prs.slide_height).inches:.2f} in)')
print()

want = [int(a) for a in sys.argv[1:]] if len(sys.argv) > 1 else list(range(1, len(prs.slides) + 1))

for i, slide in enumerate(prs.slides, 1):
    if i not in want:
        continue
    print(f'===== SLIDE {i} =====')
    print('layout:', slide.slide_layout.name)
    for sh in slide.shapes:
        kind = sh.shape_type
        pos = (f'L{Emu(sh.left).inches:.2f} T{Emu(sh.top).inches:.2f} '
               f'W{Emu(sh.width).inches:.2f} H{Emu(sh.height).inches:.2f}') if sh.left is not None else 'no-pos'
        print(f'  [{kind}] name="{sh.name}" {pos}')
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                t = ''.join(r.text for r in p.runs)
                if t.strip():
                    sizes = {r.font.size.pt for r in p.runs if r.font.size}
                    bolds = {r.font.bold for r in p.runs}
                    cols = set()
                    for r in p.runs:
                        try:
                            if r.font.color and r.font.color.rgb:
                                cols.add(str(r.font.color.rgb))
                        except Exception:
                            pass
                    print(f'      L{p.level} | {t.strip()[:150]}'
                          + (f'   <sz{sorted(sizes)} b{bolds} c{sorted(cols)}>' if (sizes or cols) else ''))
        if sh.has_table:
            tb = sh.table
            print(f'      TABLE {len(tb.rows)}x{len(tb.columns)}')
            for r_i, row in enumerate(tb.rows):
                cells = [c.text.replace('\n', ' / ').strip()[:38] for c in row.cells]
                print('        |', ' | '.join(cells))
        if sh.shape_type == 13:
            print('      PICTURE')
    print()
