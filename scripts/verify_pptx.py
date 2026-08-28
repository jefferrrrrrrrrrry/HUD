#!/usr/bin/env python3
"""PPTX 导出结果校验器（与源 deck 逐页比对）。

用法：
    python3 scripts/verify_pptx.py <deck.html> <deck.pptx>
    python3 scripts/verify_pptx.py --all

校验项：
  [0] 时效性：pptx 的修改时间不早于源 deck（防止拿旧 pptx 上台）
  [1] 页数与 deck 的 section 数一致
  [2] 文本覆盖率：deck 每页可见文本按 20 字分块，逐块检查是否出现在该页 PPTX 文本中
  [3] 越界：形状是否落在 1280 × 720 版心之内
  [4] 估算溢出：按中西文混排度量估算文本框所需高度，超出框高 118% 者列出
  [5] 表格行高合计不超过表格外框高度
  [6] 可编辑性：页面尺寸、字体统一、原生表格数、图片数不超源 deck、无空文本框
"""
from __future__ import annotations

import html as htmlmod
import math
import re
import sys
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]

DECKS = [
    ("研究汇报_2026_08.html", "研究汇报_2026_08.pptx"),
    ("文献综述_幻灯片.html", "文献综述_幻灯片.pptx"),
    ("盲区框架/研究汇报_双情境框架.html", "盲区框架/研究汇报_双情境框架.pptx"),
    ("盲区框架/开题报告_双情境框架.html", "盲区框架/开题报告_双情境框架.pptx"),
    ("盲区框架/文献综述_盲区情境.html", "盲区框架/文献综述_盲区情境.pptx"),
]

PX = 9525
W, H = 1280, 720


DECOR = r"[\s•▲▼▶●◆]+"


def norm(t: str) -> str:
    """去空白与装饰字符，使 HTML 与 PPTX 两侧口径一致。"""
    return re.sub(DECOR, "", t)


def page_texts(html: Path) -> list[str]:
    s = html.read_text(encoding="utf-8")
    secs = re.findall(r'<section class="slide.*?</section>', s, re.S)
    out = []
    for sec in secs:
        t = re.sub(r"<(script|style)\b.*?</\1>", " ", sec, flags=re.S)
        t = re.sub(r"<[^>]+>", "", t)
        t = htmlmod.unescape(t)
        out.append(norm(t))
    return out


def cjk_w(ch: str, fs: float) -> float:
    o = ord(ch)
    if o < 0x2E80 and o not in (0x2018, 0x2019, 0x201C, 0x201D):
        if ch == " ":
            return fs * 0.27
        return fs * (0.52 if ch.isalnum() else 0.33)
    return fs


def est_h(tf, box_w: float) -> float:
    total = 0.0
    for p in tf.paragraphs:
        sizes = [(r.font.size.pt if r.font.size else 10) for r in p.runs] or [10]
        lsp = p.line_spacing if isinstance(p.line_spacing, float) else 1.4
        w = sum(cjk_w(c, r.font.size.pt if r.font.size else sizes[0])
                for r in p.runs for c in r.text)
        lines = max(1, math.ceil(w / max(20.0, box_w * 0.75) - 0.02))
        # 混排段落按各 run 字号加权，避免以最大字号乘全部行数
        fs = sum(sizes) / len(sizes) if lines == 1 else max(sizes)
        total += lines * fs * lsp
    return total / 0.75  # pt → px


def check(html_rel: str, pptx_rel: str) -> int:
    html, pptx = ROOT / html_rel, ROOT / pptx_rel
    if not pptx.exists():
        print(f"  ✗ 缺失 {pptx_rel}")
        return 1
    bad = 0
    if pptx.stat().st_mtime < html.stat().st_mtime - 1:
        print(f"  [0] 时效性　✗ {pptx.name} 早于 {html.name}，须重跑 html2pptx.py")
        bad += 1
    else:
        print("  [0] 时效性　✓ pptx 不早于源 deck")
    hp = page_texts(html)
    prs = Presentation(str(pptx))
    slides = list(prs.slides)
    print(f"── {html.name} → {pptx.name}")
    print(f"  [1] 页数　deck {len(hp)}　pptx {len(slides)}"
          f"　{'✓' if len(hp) == len(slides) else '✗'}")
    if len(hp) != len(slides):
        return 1

    miss_pages, oob, over, tbad = [], [], [], []
    for i, (ht, sl) in enumerate(zip(hp, slides), 1):
        buf = []
        for sh in sl.shapes:
            if sh.has_text_frame:
                buf.append(sh.text_frame.text)
            if getattr(sh, "has_table", False) and sh.has_table:
                for row in sh.table.rows:
                    for c in row.cells:
                        buf.append(c.text)
        pt = norm("".join(buf))
        chunks = [ht[j:j + 20] for j in range(0, len(ht), 20)]
        lost = [c for c in chunks if len(c) >= 8 and c not in pt]
        if lost:
            miss_pages.append((i, len(lost), len(chunks), lost[0][:20]))

        for sh in sl.shapes:
            x, y = sh.left / PX, sh.top / PX
            w, h = sh.width / PX, sh.height / PX
            if x < -2 or y < -2 or x + w > W + 2 or y + h > H + 2:
                oob.append(f"p{i:02d} {sh.shape_type} ({x:.0f},{y:.0f},{w:.0f}×{h:.0f})")
            if sh.has_text_frame and sh.text_frame.text.strip():
                need = est_h(sh.text_frame, w)
                if need > h * 1.18 + 3:
                    over.append(f"p{i:02d} 需 {need:.0f}px / 框 {h:.0f}px　"
                                f"{sh.text_frame.text[:22]}")
            if getattr(sh, "has_table", False) and sh.has_table:
                rh = sum(r.height for r in sh.table.rows) / PX
                if rh > h + 2:
                    tbad.append(f"p{i:02d} 行高合计 {rh:.0f}px > 框 {h:.0f}px")

    if miss_pages:
        print(f"  [2] 文本覆盖　✗ {len(miss_pages)} 页有缺失")
        for i, n, tot, sample in miss_pages[:8]:
            print(f"      p{i:02d} 缺 {n}/{tot} 块　例：{sample}")
        bad += len(miss_pages)
    else:
        print("  [2] 文本覆盖　✓ 全部页逐块可查")

    print(f"  [3] 越界　{'✓ 无' if not oob else '✗ ' + str(len(oob))}")
    for o in oob[:6]:
        print("      ", o)
    bad += len(oob)

    print(f"  [4] 估算溢出　{'✓ 无' if not over else '⚠ ' + str(len(over))}")
    for o in over[:8]:
        print("      ", o)

    print(f"  [5] 表格行高　{'✓ 正常' if not tbad else '⚠ ' + str(len(tbad))}")
    for t in tbad[:6]:
        print("      ", t)

    # [6] 可编辑性
    src = html.read_text(encoding="utf-8")
    n_tab_src = len(re.findall(r"<table", src))
    n_img_src = len(re.findall(r"<img", src))
    n_tab = n_img = n_empty = 0
    fonts = set()
    for sl in slides:
        for sh in sl.shapes:
            if getattr(sh, "has_table", False) and sh.has_table:
                n_tab += 1
                for row in sh.table.rows:
                    for c in row.cells:
                        for para in c.text_frame.paragraphs:
                            for r in para.runs:
                                fonts.add(r.font.name)
            if sh.shape_type is not None and "PICTURE" in str(sh.shape_type):
                n_img += 1
            if sh.has_text_frame:
                if not sh.text_frame.text.strip() and "TEXT_BOX" in str(sh.shape_type):
                    n_empty += 1
                for para in sh.text_frame.paragraphs:
                    for r in para.runs:
                        fonts.add(r.font.name)
    size_ok = (abs(prs.slide_width / 914400 - 13.3333) < 0.01
               and abs(prs.slide_height / 914400 - 7.5) < 0.01)
    fonts.discard(None)
    bad6 = []
    if not size_ok:
        bad6.append(f"页面尺寸 {prs.slide_width / 914400:.2f}×{prs.slide_height / 914400:.2f} in")
    if fonts - {"Microsoft YaHei"}:
        bad6.append(f"字体不统一：{sorted(fonts)}")
    if n_tab != n_tab_src:
        bad6.append(f"原生表格 {n_tab} ≠ 源 deck {n_tab_src}")
    if n_img > n_img_src:
        bad6.append(f"图片 {n_img} > 源 deck {n_img_src}（正文疑被图片化）")
    if n_empty:
        bad6.append(f"空文本框 {n_empty} 个")
    if bad6:
        print(f"  [6] 可编辑性　✗ {len(bad6)} 项")
        for b in bad6:
            print("      ", b)
        bad += len(bad6)
    else:
        print(f"  [6] 可编辑性　✓ 16:9、字体统一、原生表格 {n_tab} 个、图片 {n_img} 个、无空文本框")
    return bad


def main() -> int:
    args = sys.argv[1:]
    pairs = DECKS if (not args or args[0] == "--all") else [(args[0], args[1])]
    bad = sum(check(a, b) for a, b in pairs)
    print("\n" + ("✓ 全部校验通过" if bad == 0 else f"✗ {bad} 处不通过"))
    return 0 if bad == 0 else 1


if __name__ == "__main__":
    sys.exit(main())
