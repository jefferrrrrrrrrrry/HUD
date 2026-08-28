#!/usr/bin/env python3
"""把讲述用 HTML deck 转为可编辑 PPTX。

用法：
    python3 scripts/html2pptx.py <deck.html> [out.pptx]
    python3 scripts/html2pptx.py --all            # 批量导出全部 deck

原理：以 Playwright 载入 deck，逐页读取每个渲染块的真实几何（相对 .slide 的
px 坐标）与计算样式，再在 PPTX 中以原生文本框、原生表格与矢量矩形复刻。
1280 × 720 px 的版心与 13.333 × 7.5 in 的 16:9 页面一一对应（1 px = 0.75 pt），
故字号与行距按 px × 0.75 换算即可保持版式比例。

产物为**纯原生形状**：无位图、无文本框图片化，全部文字与表格可直接编辑。
"""
from __future__ import annotations

import asyncio
import json
import sys
from pathlib import Path

from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parents[1]

DECKS = [
    "研究汇报_2026_08.html",
    "文献综述_幻灯片.html",
    "盲区框架/研究汇报_双情境框架.html",
    "盲区框架/开题报告_双情境框架.html",
    "盲区框架/文献综述_盲区情境.html",
]

PX = 9525  # 1 px = 1/96 in = 9525 EMU
FONT = "Microsoft YaHei"

# ── 浏览器侧：抽取每页的渲染块 ──────────────────────────────────────────
EXTRACT = r"""
() => {
  const INLINE = new Set(['SPAN','B','I','EM','STRONG','SUB','SUP','A','CODE','SMALL','BR']);
  const slide = document.querySelector('.slide.on');
  const sr = slide.getBoundingClientRect();
  const items = [];

  const num = v => parseFloat(v) || 0;
  const rect = el => {
    const r = el.getBoundingClientRect();
    return {x: r.left - sr.left, y: r.top - sr.top, w: r.width, h: r.height};
  };
  const rgb = c => {
    const m = /rgba?\(([\d.]+),\s*([\d.]+),\s*([\d.]+)(?:,\s*([\d.]+))?\)/.exec(c || '');
    if (!m) return null;
    const a = m[4] === undefined ? 1 : parseFloat(m[4]);
    if (a < 0.06) return null;
    return [Math.round(+m[1]), Math.round(+m[2]), Math.round(+m[3])];
  };
  const isWhite = c => c && c[0] > 250 && c[1] > 250 && c[2] > 250;

  const pseudo = el => {
    const cs = getComputedStyle(el, '::before');
    const m = /^["'](.+)["']$/.exec(cs.content || '');
    if (!m) return null;
    return {t: m[1], sz: num(cs.fontSize), c: rgb(cs.color),
            blk: cs.display.startsWith('block')};
  };

  const runsOf = (el, deferred) => {
    const st = getComputedStyle(el);
    const base = num(st.fontSize);
    const runs = [];
    const pb = pseudo(el);
    if (pb) {
      runs.push({t: pb.t, sz: pb.sz, c: pb.c});
      if (pb.blk) runs.push({t: '\n', br: true});
    }
    const walk = (node, sty) => {
      node.childNodes.forEach(ch => {
        if (ch.nodeType === 3) {
          const t = ch.textContent.replace(/\s+/g, ' ');
          if (t.trim() === '' && runs.length === 0) return;
          if (t) runs.push(Object.assign({t}, sty));
        } else if (ch.nodeType === 1) {
          if (ch.tagName === 'BR') { runs.push({t: '\n', br: true}); return; }
          if (!INLINE.has(ch.tagName)) return;
          const cs = getComputedStyle(ch);
          const bg = rgb(cs.backgroundColor);
          if (bg && !isWhite(bg)) { deferred.push(ch); return; }
          const blk = cs.display.startsWith('block');
          if (blk && runs.length) runs.push({t: '\n', br: true});
          walk(ch, {
            b: num(cs.fontWeight) >= 600 || sty.b,
            i: cs.fontStyle === 'italic' || sty.i,
            c: rgb(cs.color) || sty.c,
            sz: num(cs.fontSize) || sty.sz,
            lh: num(cs.lineHeight) || num(cs.fontSize) * 1.25 || sty.lh,
            sub: sty.sub || ch.tagName === 'SUB',
            sup: sty.sup || ch.tagName === 'SUP',
          });
          if (blk) runs.push({t: '\n', br: true});
        }
      });
    };
    walk(el, {b: num(st.fontWeight) >= 600, i: st.fontStyle === 'italic',
              c: rgb(st.color), sz: base,
              lh: num(st.lineHeight) || base * 1.25, sub: false, sup: false});
    return runs;
  };

  const decor = el => {
    const st = getComputedStyle(el);
    const bg = rgb(st.backgroundColor);
    const bw = ['Top', 'Right', 'Bottom', 'Left'].map(s => num(st['border' + s + 'Width']));
    const bc = ['Top', 'Right', 'Bottom', 'Left'].map(s => rgb(st['border' + s + 'Color']));
    const full = bw[0] > 0 && bw[1] > 0 && bw[2] > 0 && bw[3] > 0;
    const out = {};
    if (bg && !(isWhite(bg) && !full)) out.bg = bg;
    if (full) { out.line = bc[0]; out.lw = bw[0]; }
    if (!full && bw[3] >= 2.5) { out.accentL = bc[3]; out.accentLW = bw[3]; }
    if (!full && bw[2] >= 0.5 && bw[0] === 0) { out.ruleB = bc[2]; out.ruleBW = bw[2]; }
    const r = num(st.borderTopLeftRadius);
    if (r > 0) out.r = r;
    return (out.bg || out.line || out.accentL || out.ruleB) ? out : null;
  };

  const tableOf = el => {
    const rows = [];
    el.querySelectorAll('tr').forEach(tr => {
      const cells = [];
      tr.querySelectorAll('th,td').forEach(td => {
        const cs = getComputedStyle(td);
        const d = [];
        cells.push({
          runs: runsOf(td, d), rect: rect(td), head: td.tagName === 'TH',
          bg: rgb(cs.backgroundColor), align: cs.textAlign,
          pad: [num(cs.paddingTop), num(cs.paddingRight),
                num(cs.paddingBottom), num(cs.paddingLeft)],
          span: td.colSpan || 1,
        });
      });
      if (cells.length) rows.push({rect: rect(tr), cells});
    });
    return rows;
  };

  const walk = el => {
    if (el.tagName === 'TABLE') {
      const st = getComputedStyle(el);
      items.push({kind: 'table', rect: rect(el), rows: tableOf(el),
                  fs: num(st.fontSize)});
      return;
    }
    if (el.tagName === 'IMG') {
      items.push({kind: 'img', rect: rect(el), src: el.getAttribute('src')});
      return;
    }
    const d = decor(el);
    if (d) items.push(Object.assign({kind: 'rect', rect: rect(el)}, d));

    const st0 = getComputedStyle(el);
    const df = st0.display.includes('flex') || st0.display.includes('grid');
    const kids = Array.from(el.children);
    const leaf = !(df && kids.length >= 2) && kids.every(k => INLINE.has(k.tagName));
    if (leaf) {
      const deferred = [];
      const runs = runsOf(el, deferred);
      if (runs.some(r => r.t.trim())) {
        const st = getComputedStyle(el);
        const li = el.tagName === 'LI';
        const r0 = rect(el);
        if (li) { r0.x -= 15; r0.w += 15; }
        let align = st.textAlign, vc = false;
        if (df) {
          const colDir = (st.flexDirection || '').startsWith('column');
          if (colDir) {
            if (st.alignItems === 'center') align = 'center';
            vc = st.justifyContent === 'center';
          } else {
            vc = st.alignItems === 'center';
            if (st.justifyContent === 'center') align = 'center';
          }
        }
        items.push({
          kind: 'text', rect: r0, runs,
          fs: num(st.fontSize), lh: num(st.lineHeight) || num(st.fontSize) * 1.25,
          align: align, bullet: li, vc: vc,
          pad: [num(st.paddingTop), num(st.paddingRight),
                num(st.paddingBottom), num(st.paddingLeft)],
        });
      }
      deferred.forEach(walk);
    } else {
      // flex/grid 容器的裸文本节点会形成匿名项，须以 Range 量出其框后单独输出
      Array.from(el.childNodes).forEach(n => {
        if (n.nodeType === 1) { walk(n); return; }
        if (n.nodeType !== 3 || !n.textContent.trim()) return;
        const rng = document.createRange();
        rng.selectNode(n);
        const rr = rng.getBoundingClientRect();
        if (!rr.width || !rr.height) return;
        const st = getComputedStyle(el);
        items.push({
          kind: 'text',
          rect: {x: rr.left - sr.left, y: rr.top - sr.top,
                 w: rr.width + 2, h: rr.height},
          runs: [{t: n.textContent.replace(/\s+/g, ' '),
                  b: num(st.fontWeight) >= 600, i: st.fontStyle === 'italic',
                  c: rgb(st.color), sz: num(st.fontSize)}],
          fs: num(st.fontSize),
          lh: num(st.lineHeight) || num(st.fontSize) * 1.4,
          align: 'center', bullet: false, vc: true, pad: [0, 0, 0, 0],
        });
      });
    }
  };

  Array.from(slide.children).forEach(walk);
  return {w: sr.width, h: sr.height, items};
}
"""


async def grab(html: Path) -> list[dict]:
    async with async_playwright() as pw:
        br = await pw.chromium.launch()
        pg = await br.new_page(viewport={"width": 1280, "height": 720})
        await pg.goto(html.resolve().as_uri())
        await pg.wait_for_timeout(320)
        n = await pg.evaluate("() => document.querySelectorAll('.slide').length")
        pages = []
        for i in range(n):
            await pg.evaluate(
                "(i) => { const ss=[...document.querySelectorAll('.slide')];"
                "ss.forEach(s=>s.classList.remove('on')); ss[i].classList.add('on');"
                "ss[i].style.transform='none';}", i)
            await pg.wait_for_timeout(90)
            pages.append(await pg.evaluate(EXTRACT))
        await br.close()
        return pages


# ── PPTX 侧 ──────────────────────────────────────────────────────────────
def col(c) -> RGBColor:
    return RGBColor(*c)


def no_shadow(shape):
    spPr = shape._element.spPr
    for tag in ("a:effectLst", "a:effectDag"):
        for e in spPr.findall(qn(tag)):
            spPr.remove(e)
    spPr.append(spPr.makeelement(qn("a:effectLst"), {}))


def put_rect(slide, it):
    r = it["rect"]
    if r["w"] < 1 or r["h"] < 1:
        return
    if not (it.get("bg") or it.get("line") or it.get("accentL")):
        if it.get("ruleB"):
            w = max(0.75, it["ruleBW"])
            ln = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Emu(int(r["x"] * PX)),
                Emu(int((r["y"] + r["h"] - w) * PX)),
                Emu(int(r["w"] * PX)), Emu(int(w * PX)))
            ln.fill.solid()
            ln.fill.fore_color.rgb = col(it["ruleB"])
            ln.line.fill.background()
            ln.shadow.inherit = False
            no_shadow(ln)
        return
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if it.get("r", 0) >= 3 else MSO_SHAPE.RECTANGLE,
        Emu(int(r["x"] * PX)), Emu(int(r["y"] * PX)),
        Emu(int(r["w"] * PX)), Emu(int(r["h"] * PX)))
    if it.get("r", 0) >= 3:
        try:
            shape.adjustments[0] = min(0.5, it["r"] / min(r["w"], r["h"]))
        except (IndexError, ZeroDivisionError):
            pass
    if it.get("bg"):
        shape.fill.solid()
        shape.fill.fore_color.rgb = col(it["bg"])
    else:
        shape.fill.background()
    if it.get("line"):
        shape.line.color.rgb = col(it["line"])
        shape.line.width = Pt(max(0.5, it.get("lw", 1) * 0.75))
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    no_shadow(shape)
    shape.text_frame.text = ""

    if it.get("accentL"):
        w = it["accentLW"]
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(int(r["x"] * PX)), Emu(int(r["y"] * PX)),
            Emu(int(w * PX)), Emu(int(r["h"] * PX)))
        bar.fill.solid()
        bar.fill.fore_color.rgb = col(it["accentL"])
        bar.line.fill.background()
        bar.shadow.inherit = False
        no_shadow(bar)
    if it.get("ruleB"):
        w = max(0.75, it["ruleBW"])
        ln = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(int(r["x"] * PX)),
            Emu(int((r["y"] + r["h"] - w) * PX)),
            Emu(int(r["w"] * PX)), Emu(int(w * PX)))
        ln.fill.solid()
        ln.fill.fore_color.rgb = col(it["ruleB"])
        ln.line.fill.background()
        ln.shadow.inherit = False
        no_shadow(ln)


ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
         "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY,
         "start": PP_ALIGN.LEFT, "end": PP_ALIGN.RIGHT}


def fill_runs(tf, runs, fs, lh, align, bullet=False):
    paras = [[]]
    for r in runs:
        if r.get("br"):
            paras.append([])
        else:
            paras[-1].append(r)
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for chunk in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = ALIGN.get(align, PP_ALIGN.LEFT)
        # 行距取本段最大字号那一 run 自身的行高，避免沿用父级比例导致行高虚高
        if chunk:
            dom = max(chunk, key=lambda r: r.get("sz") or fs)
            pmax = dom.get("sz") or fs
            plh = dom.get("lh") or lh
        else:
            pmax, plh = fs, lh
        if plh and pmax:
            p.line_spacing = round(min(2.4, max(1.02, plh / pmax)), 3)
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        if bullet and chunk:
            b = p.add_run()
            b.text = "• "
            b.font.size = Pt(round(fs * 0.75, 1))
            b.font.name = FONT
        for r in chunk:
            run = p.add_run()
            run.text = r["t"]
            f = run.font
            f.name = FONT
            f.size = Pt(round((r.get("sz") or fs) * 0.75, 1))
            f.bold = bool(r.get("b"))
            f.italic = bool(r.get("i"))
            if r.get("c"):
                f.color.rgb = col(r["c"])
            if r.get("sub"):
                run.font._rPr.set("baseline", "-25000")
            elif r.get("sup"):
                run.font._rPr.set("baseline", "30000")


def put_text(slide, it):
    r, pad = it["rect"], it.get("pad", [0, 0, 0, 0])
    box = slide.shapes.add_textbox(  # noqa: E501
        Emu(int((r["x"] + pad[3]) * PX)), Emu(int((r["y"] + pad[0]) * PX)),
        Emu(int(max(12, r["w"] - pad[1] - pad[3] + 4) * PX)),
        Emu(int(max(10, r["h"] - pad[0] - pad[2] + 3) * PX)))
    fill_runs(box.text_frame, it["runs"], it["fs"], it["lh"],
              it.get("align", "left"), it.get("bullet", False))
    if it.get("vc"):
        box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def put_table(slide, it):
    rows = it["rows"]
    if not rows:
        return
    ncol = max(sum(c["span"] for c in row["cells"]) for row in rows)
    r = it["rect"]
    gf = slide.shapes.add_table(
        len(rows), ncol, Emu(int(r["x"] * PX)), Emu(int(r["y"] * PX)),
        Emu(int(r["w"] * PX)), Emu(int(r["h"] * PX)))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False

    widest = max(rows, key=lambda x: len(x["cells"]))
    if len(widest["cells"]) == ncol:
        for j, c in enumerate(widest["cells"]):
            tbl.columns[j].width = Emu(int(c["rect"]["w"] * PX))
    for i, row in enumerate(rows):
        tbl.rows[i].height = Emu(int(row["rect"]["h"] * PX))
        j = 0
        for c in row["cells"]:
            if j >= ncol:
                break
            cell = tbl.cell(i, j)
            if c["span"] > 1 and j + c["span"] - 1 < ncol:
                cell.merge(tbl.cell(i, j + c["span"] - 1))
            pad = c["pad"]
            cell.margin_left, cell.margin_right = Emu(int(pad[3] * PX)), Emu(int(pad[1] * PX))
            cell.margin_top, cell.margin_bottom = Emu(int(pad[0] * PX)), Emu(int(pad[2] * PX))
            cell.vertical_anchor = MSO_ANCHOR.TOP
            if c["bg"]:
                cell.fill.solid()
                cell.fill.fore_color.rgb = col(c["bg"])
            else:
                cell.fill.background()
            fs = (c["runs"][0].get("sz") if c["runs"] else None) or it["fs"]
            fill_runs(cell.text_frame, c["runs"], fs, fs * 1.55, c["align"])
            j += c["span"]


def put_img(slide, it, base: Path):
    src = (base / it["src"]).resolve()
    if not src.exists():
        return
    r = it["rect"]
    slide.shapes.add_picture(str(src), Emu(int(r["x"] * PX)), Emu(int(r["y"] * PX)),
                             Emu(int(r["w"] * PX)), Emu(int(r["h"] * PX)))


def build(pages: list[dict], out: Path, base: Path) -> None:
    prs = Presentation()
    prs.slide_width = Emu(1280 * PX)
    prs.slide_height = Emu(720 * PX)
    blank = prs.slide_layouts[6]
    for pg in pages:
        slide = prs.slides.add_slide(blank)
        for it in pg["items"]:
            k = it["kind"]
            if k == "rect":
                put_rect(slide, it)
            elif k == "text":
                put_text(slide, it)
            elif k == "table":
                put_table(slide, it)
            elif k == "img":
                put_img(slide, it, base)
    prs.save(str(out))


def convert(rel: str, out: str | None = None) -> Path:
    html = ROOT / rel
    dst = Path(out) if out else html.with_suffix(".pptx")
    if not dst.is_absolute():
        dst = ROOT / dst
    pages = asyncio.run(grab(html))
    build(pages, dst, html.parent)
    n_shape = sum(len(p["items"]) for p in pages)
    try:
        shown = dst.relative_to(ROOT)
    except ValueError:
        shown = dst
    print(f"✓ {shown}　{len(pages)} 页　{n_shape} 个可编辑对象")
    return dst


def main() -> int:
    args = sys.argv[1:]
    if not args:
        print(__doc__)
        return 1
    if args[0] == "--all":
        for d in DECKS:
            convert(d)
        return 0
    convert(args[0], args[1] if len(args) > 1 else None)
    return 0


if __name__ == "__main__":
    sys.exit(main())
