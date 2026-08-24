#!/usr/bin/env python3
"""Generate the open-defense PPT for HUD/AR-HUD pedestrian warning research.

Uses python-pptx with a clean academic style.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
import pathlib

OUTPUT = "/home/gezhuocheng/HUD/PPT_开题报告_HUD_AR-HUD时空设计规范.pptx"

# Color palette (similar to 周颖 PPT style - blue gradient)
COLOR_PRIMARY = RGBColor(0x1F, 0x3A, 0x68)        # deep blue
COLOR_SECONDARY = RGBColor(0x4A, 0x6F, 0xA5)      # medium blue
COLOR_ACCENT = RGBColor(0xE0, 0x52, 0x39)         # red accent
COLOR_LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)       # very light blue
COLOR_TEXT = RGBColor(0x21, 0x21, 0x21)           # near black
COLOR_TEXT_SOFT = RGBColor(0x55, 0x55, 0x55)
COLOR_HIGHLIGHT = RGBColor(0xFD, 0xC7, 0x00)      # yellow highlight

# Build new presentation 16:9 (default), then we'll resize
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ---- helpers ----

def add_blank_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    return slide

def add_filled_rect(slide, left, top, width, height, fill_color, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    sh.line.fill.background() if line is None else None
    if line is None:
        sh.line.fill.background()
    return sh

def add_text_box(slide, left, top, width, height, text, font_size=18, 
                 bold=False, color=COLOR_TEXT, align=PP_ALIGN.LEFT, 
                 font_name="Microsoft YaHei", vertical=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = vertical
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(font_size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font_name
    return tb

def add_slide_title_bar(slide, title_text, subtitle_text=""):
    """Add a colored title bar at top + title text."""
    # Top color bar
    add_filled_rect(slide, 0, 0, prs.slide_width, Inches(1.0), COLOR_PRIMARY)
    # Title text
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(11), Inches(0.7),
                 title_text, font_size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                 vertical=MSO_ANCHOR.MIDDLE)
    if subtitle_text:
        add_text_box(slide, Inches(0.5), Inches(0.7), Inches(11), Inches(0.3),
                     subtitle_text, font_size=14, color=RGBColor(0xCC, 0xDD, 0xEE),
                     vertical=MSO_ANCHOR.TOP)

def add_page_number(slide, num, total=34):
    add_text_box(slide, Inches(11.8), Inches(7.0), Inches(1.4), Inches(0.3),
                 f"{num} / {total}", font_size=10, color=COLOR_TEXT_SOFT,
                 align=PP_ALIGN.RIGHT)

def add_footer_line(slide):
    add_filled_rect(slide, 0, Inches(7.25), prs.slide_width, Emu(38100), COLOR_SECONDARY)

# ---- Slide builders ----

def make_cover_slide():
    slide = add_blank_slide()
    # Big background
    add_filled_rect(slide, 0, 0, prs.slide_width, prs.slide_height, COLOR_PRIMARY)
    # Accent stripe
    add_filled_rect(slide, 0, Inches(3.4), prs.slide_width, Inches(0.05), COLOR_ACCENT)
    # Title
    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.2),
                 "HUD/AR-HUD 行人碰撞预警的", font_size=42, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.0),
                 "时间-空间元素设计规范研究", font_size=40, bold=True,
                 color=COLOR_HIGHLIGHT, align=PP_ALIGN.CENTER)
    # English subtitle
    add_text_box(slide, Inches(0.8), Inches(3.7), Inches(11.7), Inches(0.6),
                 "Design Specification for Temporal-Spatial Elements in HUD/AR-HUD Pedestrian Collision Warning",
                 font_size=18, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER,
                 font_name="Calibri")
    # Author info
    add_text_box(slide, Inches(2), Inches(5.0), Inches(9.3), Inches(0.5),
                 "答辩人：[姓名]", font_size=22, color=RGBColor(0xFF, 0xFF, 0xFF),
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), Inches(5.5), Inches(9.3), Inches(0.5),
                 "指导老师：[导师]", font_size=20, color=RGBColor(0xCC, 0xDD, 0xEE),
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), Inches(6.5), Inches(9.3), Inches(0.5),
                 "2026年[月]", font_size=18, color=RGBColor(0xCC, 0xDD, 0xEE),
                 align=PP_ALIGN.CENTER)
    return slide

def make_toc_slide(current_section=None, total=5):
    slide = add_blank_slide()
    add_slide_title_bar(slide, "目录", "Contents")
    sections = [
        ("1", "研究背景与意义"),
        ("2", "文献综述"),
        ("3", "研究问题"),
        ("4", "研究方案"),
        ("5", "研究工作计划与进度安排"),
    ]
    y_start = Inches(1.5)
    for i, (num, name) in enumerate(sections):
        y = y_start + Inches(0.9 * i)
        # number circle
        if current_section == num:
            color = COLOR_ACCENT
            text_color = RGBColor(0xFF, 0xFF, 0xFF)
        else:
            color = COLOR_LIGHT_BG
            text_color = COLOR_PRIMARY
        # Circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.5), y, Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        # Number inside circle
        add_text_box(slide, Inches(2.5), y, Inches(0.7), Inches(0.7),
                     num, font_size=24, bold=True, color=text_color,
                     align=PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE)
        # Name to the right
        name_color = COLOR_ACCENT if current_section == num else COLOR_TEXT
        add_text_box(slide, Inches(3.5), y, Inches(8), Inches(0.7),
                     name, font_size=22, bold=(current_section == num),
                     color=name_color, vertical=MSO_ANCHOR.MIDDLE)
    return slide

def make_section_title_slide(section_num, title, en_title, total=34, page=None):
    slide = add_blank_slide()
    add_filled_rect(slide, 0, 0, prs.slide_width, prs.slide_height, COLOR_PRIMARY)
    # Large section number
    add_text_box(slide, Inches(1.5), Inches(2.2), Inches(2.5), Inches(2.5),
                 section_num, font_size=180, bold=True,
                 color=COLOR_HIGHLIGHT, align=PP_ALIGN.CENTER,
                 vertical=MSO_ANCHOR.MIDDLE, font_name="Calibri")
    # Decoration line
    add_filled_rect(slide, Inches(5), Inches(3.2), Inches(0.05), Inches(1.0), COLOR_ACCENT)
    # Title
    add_text_box(slide, Inches(5.5), Inches(3.0), Inches(7.5), Inches(0.8),
                 title, font_size=36, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF))
    # English subtitle
    add_text_box(slide, Inches(5.5), Inches(3.8), Inches(7.5), Inches(0.5),
                 en_title, font_size=16,
                 color=RGBColor(0xCC, 0xDD, 0xEE), font_name="Calibri")
    if page:
        add_page_number(slide, page, total)
    return slide

def make_content_slide(title, bullets, page=None, total=34, subtitle=None):
    """Make a content slide with title bar + bullets."""
    slide = add_blank_slide()
    add_slide_title_bar(slide, title, subtitle or "")
    add_footer_line(slide)
    # Bullets
    y = Inches(1.3)
    bullet_text_lines = []
    for b in bullets:
        if isinstance(b, tuple):  # (level, text)
            level, text = b
            prefix = "  " * level + ("· " if level == 0 else ("- " if level == 1 else "  · "))
            bullet_text_lines.append(prefix + text)
        else:
            bullet_text_lines.append("· " + b)
    
    text = "\n".join(bullet_text_lines)
    add_text_box(slide, Inches(0.6), y, Inches(12.2), Inches(5.7),
                 text, font_size=16, color=COLOR_TEXT)
    if page:
        add_page_number(slide, page, total)
    return slide

def make_richtext_slide(title, paragraphs, page=None, total=34, subtitle=None):
    """Make a slide with mixed-formatting paragraphs.
    
    paragraphs = list of dicts: {"text": str, "size": int, "bold": bool, "color": RGBColor, "indent": int}
    """
    slide = add_blank_slide()
    add_slide_title_bar(slide, title, subtitle or "")
    add_footer_line(slide)
    
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.2), Inches(5.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    
    for i, para in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Set indent
        if "indent" in para:
            p.level = para["indent"]
        r = p.add_run()
        r.text = para["text"]
        r.font.size = Pt(para.get("size", 16))
        r.font.bold = para.get("bold", False)
        r.font.color.rgb = para.get("color", COLOR_TEXT)
        r.font.name = para.get("font", "Microsoft YaHei")
    
    if page:
        add_page_number(slide, page, total)
    return slide

# ============ Build all slides ============

# Track total slides for page numbering
SLIDES = []

# Slide 1: Cover
SLIDES.append(("cover", make_cover_slide))
# Slide 2: TOC
SLIDES.append(("toc", lambda: make_toc_slide()))

# Section 1: 研究背景与意义 (slides 3-6)
SLIDES.append(("section1", lambda: make_section_title_slide("1", "研究背景与意义", "Research Background and Significance")))
SLIDES.append(("bg_intro", None))  # slide 4
SLIDES.append(("bg_problems", None))  # slide 5
SLIDES.append(("bg_significance", None))  # slide 6

# Slide 7: TOC highlight section 2
SLIDES.append(("toc2", lambda: make_toc_slide("2")))

# Section 2: 文献综述 (slides 8-22)
SLIDES.append(("section2", lambda: make_section_title_slide("2", "文献综述", "Literature Review")))
SLIDES.append(("lit_overview", None))  # slide 9
SLIDES.append(("lit_concepts_t", None))  # slide 10
SLIDES.append(("lit_concepts_s", None))  # slide 11
SLIDES.append(("lit_t_onset", None))  # slide 12
SLIDES.append(("lit_t_duration", None))  # slide 13
SLIDES.append(("lit_s_color", None))  # slide 14
SLIDES.append(("lit_s_shape", None))  # slide 15
SLIDES.append(("lit_s_layout", None))  # slide 16
SLIDES.append(("lit_individual", None))  # slide 17
SLIDES.append(("lit_cognitive", None))  # slide 18
SLIDES.append(("lit_gaps", None))  # slide 19
SLIDES.append(("lit_theory", None))  # slide 20

# Slide 21: TOC highlight section 3
SLIDES.append(("toc3", lambda: make_toc_slide("3")))

# Section 3: 研究问题 (slides 22-24)
SLIDES.append(("section3", lambda: make_section_title_slide("3", "研究问题", "Research Questions")))
SLIDES.append(("rq_summary", None))  # slide 23
SLIDES.append(("rq_framework", None))  # slide 24

# Slide 25: TOC highlight section 4
SLIDES.append(("toc4", lambda: make_toc_slide("4")))

# Section 4: 研究方案 (slides 26-31)
SLIDES.append(("section4", lambda: make_section_title_slide("4", "研究方案", "Research Plan")))
SLIDES.append(("plan_overview", None))  # slide 27
SLIDES.append(("exp1", None))  # slide 28
SLIDES.append(("exp2a", None))  # slide 29
SLIDES.append(("exp2b", None))  # slide 30
SLIDES.append(("exp3", None))  # slide 31
SLIDES.append(("plan_route", None))  # slide 32
SLIDES.append(("plan_summary_table", None))  # slide 33

# Slide 34: TOC highlight section 5
SLIDES.append(("toc5", lambda: make_toc_slide("5")))

# Section 5: 工作计划 (slides 35-36)
SLIDES.append(("section5", lambda: make_section_title_slide("5", "研究工作计划与进度安排", "Work Plan")))
SLIDES.append(("plan_gantt", None))  # slide 36
SLIDES.append(("plan_outputs", None))  # slide 37

# Thanks
SLIDES.append(("thanks", None))  # slide 38

TOTAL = len(SLIDES)
print(f"Planned {TOTAL} slides")

# Make slide 1 (cover) and 2 (TOC) - done via lambda
# Reset and build properly

# Clear and rebuild
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide 1: Cover
make_cover_slide()
# Slide 2: TOC
make_toc_slide()
# Slide 3: Section 1 title
make_section_title_slide("1", "研究背景与意义", "Research Background and Significance", total=TOTAL, page=3)

# Slide 4: Background intro
make_richtext_slide(
    "1.1 研究背景",
    [
        {"text": "HUD/AR-HUD 技术正快速进入量产车型", "size": 22, "bold": True, "color": COLOR_PRIMARY},
        {"text": " ", "size": 8},
        {"text": "·  据 IDC 预测，2030 年全球 AR-HUD 装机量将达约百万套量级", "size": 17},
        {"text": "·  2024 年蔚来、奔驰、奥迪等品牌已大规模量产 AR-HUD 系统", "size": 17},
        {"text": "·  行人事故是道路致死的第一大成因（WHO, 2023：约 35%）", "size": 17, "color": COLOR_ACCENT},
        {"text": " ", "size": 8},
        {"text": "AR-HUD 为行人碰撞预警提供新的可能性", "size": 22, "bold": True, "color": COLOR_PRIMARY},
        {"text": " ", "size": 8},
        {"text": "·  共形增强（contact-analog）：将虚拟图形与真实行人在视野中精准对齐", "size": 17},
        {"text": "·  无需视线下移，减少注意力切换成本", "size": 17},
        {"text": "·  代表性研究：Kim 等（2018）、Wu 等（2024）", "size": 15, "color": COLOR_TEXT_SOFT},
    ],
    page=4, total=TOTAL
)

# Slide 5: Problems
make_richtext_slide(
    "1.2 研究问题的现实困境",
    [
        {"text": "AR-HUD 行人预警的『时空设计规范』尚未确立", "size": 22, "bold": True, "color": COLOR_PRIMARY},
        {"text": " ", "size": 8},
        {"text": "时间维度争议：", "size": 18, "bold": True, "color": COLOR_ACCENT},
        {"text": "·  何时出现？TTC=2.5s vs 5.0s（Kim et al., 2018）vs 100m 前提示（Zhang et al., 2024）", "size": 15},
        {"text": "·  持续多久？显示到危险解除 vs 固定 3-15s（Ma et al., 2021）", "size": 15},
        {"text": "·  是否分级？两级 vs 三级 vs 不分级", "size": 15},
        {"text": " ", "size": 8},
        {"text": "空间维度争议：", "size": 18, "bold": True, "color": COLOR_ACCENT},
        {"text": "·  锁定方式：屏幕固定 vs 行人共形 vs 路面共形（Wu et al., 2024）", "size": 15},
        {"text": "·  颜色编码：单色红 vs 红黄绿渐变（Ma et al., 2024）vs 优先级编码（Chen et al., 2025）", "size": 15},
        {"text": "·  动效：静态 vs 动态共形跟随 vs 闪烁（Huo & Alla, 2025）", "size": 15},
        {"text": " ", "size": 8},
        {"text": "→ 现有研究碎片化，缺乏跨时空维度的整合设计规范", "size": 18, "bold": True, "color": COLOR_PRIMARY},
    ],
    page=5, total=TOTAL
)

# Slide 6: Significance
make_richtext_slide(
    "1.3 研究意义",
    [
        {"text": "理论意义", "size": 20, "bold": True, "color": COLOR_ACCENT},
        {"text": "·  整合时间×空间维度，建立 HUD/AR-HUD 行人预警的统一设计框架", "size": 16},
        {"text": "·  验证经典视觉注意理论（Wickens 多资源理论、Mack & Rock 非注意盲视）", "size": 16},
        {"text": "    在 AR-HUD 场景下的适用性", "size": 16},
        {"text": " ", "size": 8},
        {"text": "实践意义", "size": 20, "bold": True, "color": COLOR_ACCENT},
        {"text": "·  为量产车 AR-HUD 的 HMI 设计提供量化规范（颜色 RGB / 形状 / FOV / TTC 阈值）", "size": 16},
        {"text": "·  重点服务新手驾驶员（占国内新增驾照 70%）的安全需求", "size": 16},
        {"text": " ", "size": 8},
        {"text": "应用价值", "size": 20, "bold": True, "color": COLOR_ACCENT},
        {"text": "·  输出可被 ISO 15008 / SAE J2400 等标准引用的设计参数表", "size": 16},
        {"text": "·  推动新手驾驶员培训系统的可视化辅助标准化", "size": 16},
    ],
    page=6, total=TOTAL
)

# Slide 7: TOC2
make_toc_slide("2")

# Slide 8: Section 2 title
make_section_title_slide("2", "文献综述", "Literature Review", total=TOTAL, page=8)

# Slide 9: Literature overview
make_richtext_slide(
    "2.1 文献检索方案",
    [
        {"text": "检索范围", "size": 20, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  关键词：HUD / AR-HUD / pedestrian / collision warning / contact-analog / virtual shadow", "size": 15},
        {"text": "·  时间：2008–2025", "size": 15},
        {"text": "·  数据库：Web of Science、Scopus、IEEE Xplore、OpenAlex、CNKI", "size": 15},
        {"text": "·  补充检索：科研通互助、sci-hub、Cambridge Apollo、ResearchGate", "size": 15},
        {"text": " ", "size": 8},
        {"text": "最终纳入：40 篇核心文献", "size": 20, "bold": True, "color": COLOR_ACCENT},
        {"text": " ", "size": 8},
        {"text": "·  期刊论文 31 篇（含 IEEE TVCG、IJHCI、Sensors、Sustainability、IEEE Access）", "size": 15},
        {"text": "·  会议论文 6 篇（CHI、AutomotiveUI、IEEE ITSC、ICCE 等）", "size": 15},
        {"text": "·  综述论文 3 篇（Winkler & Soleimani 2025、Kettle & Lee 2022、Skirnewskaja 2022）", "size": 15},
        {"text": " ", "size": 8},
        {"text": "按时间分布：2008–2015 (4) → 2016–2020 (12) → 2021–2025 (24)", "size": 15, "color": COLOR_TEXT_SOFT},
        {"text": "→ 研究呈爆发增长趋势，反映 AR-HUD 量产化驱动学术关注", "size": 16, "bold": True, "color": COLOR_PRIMARY},
    ],
    page=9, total=TOTAL
)

# Slide 10: Concepts - Time
make_richtext_slide(
    "2.2 核心概念界定（时间维度）",
    [
        {"text": "TTC (Time-to-Collision)  碰撞时间", "size": 18, "bold": True, "color": COLOR_ACCENT},
        {"text": "·  定义：以当前相对速度恒定计算，本车与目标到达同一位置的剩余时间", "size": 14},
        {"text": "·  来源：Hayward (1972) 在 Highway Research Record 首次系统提出", "size": 14},
        {"text": "·  公式：TTC = d / v_rel；典型阈值：2-5s 用于警告触发", "size": 14},
        {"text": " ", "size": 6},
        {"text": "PRT (Perception-Reaction Time)  感知反应时间", "size": 18, "bold": True, "color": COLOR_ACCENT},
        {"text": "·  PIEV 模型 (Hooper, 1936): Perception-Identification-Emotion-Volition", "size": 14},
        {"text": "·  典型范围：0.7-1.5s（清醒驾驶员）；分心条件下显著延长", "size": 14},
        {"text": " ", "size": 6},
        {"text": "TTFF (Time to First Fixation)  首次注视时间", "size": 18, "bold": True, "color": COLOR_ACCENT},
        {"text": "·  来源：眼动追踪研究标准指标 (Rayner, 1998)", "size": 14},
        {"text": "·  在 AR-HUD 评估中：衡量警示捕获注意的速度", "size": 14},
        {"text": " ", "size": 6},
        {"text": "Warning Onset / Duration  警告出现时机 / 持续时长", "size": 18, "bold": True, "color": COLOR_ACCENT},
        {"text": "·  Onset：通常以 TTC 阈值或距离阈值表征", "size": 14},
        {"text": "·  Duration：可分为『至危险解除』模式与『固定时长』模式", "size": 14},
    ],
    page=10, total=TOTAL
)

# Slide 11: Concepts - Space
make_richtext_slide(
    "2.3 核心概念界定（空间维度）",
    [
        {"text": "Contact-Analog / Conformal  接触类比 / 共形", "size": 18, "bold": True, "color": COLOR_ACCENT},
        {"text": "·  Tönnis et al. (2007)：虚拟图形与真实世界对象在驾驶员视野中精确对齐", "size": 14},
        {"text": "·  核心子类：行人锁定 / 路面锁定 / 世界锁定 vs 屏幕固定 (Wu et al., 2024)", "size": 14},
        {"text": " ", "size": 6},
        {"text": "FOV (Field of View)  视场角", "size": 18, "bold": True, "color": COLOR_ACCENT},
        {"text": "·  HUD 虚像可见的水平×垂直角度", "size": 14},
        {"text": "·  量产 HUD：10-20°×3-5°；AR-HUD：25-40°×7-12°", "size": 14},
        {"text": " ", "size": 6},
        {"text": "颜色与可见度参数", "size": 18, "bold": True, "color": COLOR_ACCENT},
        {"text": "·  CIE 色坐标 (x, y)：1931 年 CIE 国际照明委员会色度图", "size": 14},
        {"text": "·  亮度 Lv (cd/m²)、对比度 Weber/Michelson", "size": 14},
        {"text": "·  颜色语义编码：红=danger / 黄=caution / 绿=safe（源于 SAE J2400、交通灯）", "size": 14},
        {"text": " ", "size": 6},
        {"text": "Gaze Entropy  注视熵", "size": 18, "bold": True, "color": COLOR_ACCENT},
        {"text": "·  基于 Shannon (1948) 信息熵；量化注视分布的随机性", "size": 14},
        {"text": "·  熵 ↑ = 注意力分散；熵 ↓ = 注意力集中", "size": 14},
    ],
    page=11, total=TOTAL
)

# Slide 12: Lit - Warning Onset
make_richtext_slide(
    "2.4 文献综述（1）警告出现时机",
    [
        {"text": "固定 TTC 阈值类研究", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  Kim et al. (2018, IEEE TVCG)：2.5s（近）/5.0s（远）双距离条件 + GPS 触发", "size": 14},
        {"text": "·  Lübbe (2017, J Safety Res)：cautionary HUD @ TTC=2.5s → imminent audio-visual @ 1.8s，间隔 0.7s", "size": 14},
        {"text": "·  Huo & Alla (2025)：TTC=2.5s / 距离 34.72m 触发，flashing 动态明确", "size": 14},
        {"text": "·  Wu et al. (2024)：TTC<3s 触发；60m 处行人激活", "size": 14},
        {"text": " ", "size": 6},
        {"text": "距离阈值类研究", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  Phan et al. (2016, ITSC)：tWP = min(t(TTC=2s), t(d=16.6m))，组合阈值", "size": 14},
        {"text": "·  Zhang et al. (2024, 中文版华南理工大学学报)：100m 前语音预警 + 60m 行人激活", "size": 14},
        {"text": " ", "size": 6},
        {"text": "自适应触发类研究", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  Frémont et al. (2019)：基于头眼监测的自适应触发", "size": 14},
        {"text": "·  Doshi et al. (2008)：主动注视方向跟踪 + 预测性 HUD", "size": 14},
        {"text": " ", "size": 6},
        {"text": "→ 问题提出 1：警告出现时机的最优 TTC 阈值是多少？", "size": 16, "bold": True, "color": COLOR_ACCENT},
        {"text": "    是否随驾驶员经验 / 场景变化？", "size": 16, "color": COLOR_ACCENT},
    ],
    page=12, total=TOTAL
)

# Slide 13: Lit - Duration & Graded
make_richtext_slide(
    "2.5 文献综述（2）警告持续时长与分级",
    [
        {"text": "持续时长：两种模式", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  『至危险解除』模式（多数研究）：Kim 2018、Phan 2016、Zhang 2024", "size": 14},
        {"text": "·  『固定时长』模式：Ma (2021, IEEE Access) 3s 常规 / 10-15s 紧急", "size": 14},
        {"text": " ", "size": 6},
        {"text": "分级警告设计", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  两级（cautionary → imminent）：Lübbe 2017（间隔 0.7s）", "size": 14},
        {"text": "·  三级威胁等级：Yoon 2014（基于 TTC 计算，未量化阈值）", "size": 14},
        {"text": "·  颜色×饱和度渐变三级：Ma 2024 EID（绿/黄/红 + saturation 渐变）", "size": 14},
        {"text": "·  优先级颜色编码：Chen et al. 2025 IJHCI（多目标场景，N=45 新手）", "size": 14, "bold": True, "color": COLOR_ACCENT},
        {"text": " ", "size": 6},
        {"text": "速度自适应分级", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  Ma (2021)、Teng (2023)：FOV 与显示密度随车速分级", "size": 14},
        {"text": " ", "size": 6},
        {"text": "→ 问题提出 2：分级警告的级别数与级别间隔如何最优化？", "size": 16, "bold": True, "color": COLOR_ACCENT},
        {"text": "    单目标 vs 多目标场景下的差异？", "size": 16, "color": COLOR_ACCENT},
    ],
    page=13, total=TOTAL
)

# Slide 14: Lit - Color
make_richtext_slide(
    "2.6 文献综述（3）色彩编码",
    [
        {"text": "单色警示主流", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  红色：Kim 2018、Wu 2024 RGB(255,0,0)、Kazazi 2015", "size": 14},
        {"text": "·  黄色：Phan 2016（避免过度紧迫感）", "size": 14},
        {"text": " ", "size": 6},
        {"text": "分级配色", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  红/黄/绿 + 饱和度渐变：Ma 2024 EID 'carpet' 设计", "size": 14},
        {"text": "·  4 色 HEX 分区：Teng 2023（#2979FF / #FE0000 / #4ADE80 / #F26D21）", "size": 14},
        {"text": "·  红=紧急 + 青蓝=辅助：Ma 2021", "size": 14},
        {"text": "·  优先级颜色：Chen 2025（Hierarchical vs Equivalent vs Baseline）", "size": 14, "bold": True, "color": COLOR_ACCENT},
        {"text": " ", "size": 6},
        {"text": "色彩可见度专题", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  Zhong (2022)：7 色 × 3 描边 × 2 照度（100,000lx vs 30lx）", "size": 14},
        {"text": "·  发现：高照度下饱和度差异减小，需通过描边补偿", "size": 14},
        {"text": " ", "size": 6},
        {"text": "→ 颜色研究共识：『红=紧急』；分级配色优于单色", "size": 16, "bold": True, "color": COLOR_PRIMARY},
        {"text": "→ 空白：中国驾驶员的颜色语义偏好尚未系统研究", "size": 16, "color": COLOR_ACCENT},
    ],
    page=14, total=TOTAL
)

# Slide 15: Lit - Shape
make_richtext_slide(
    "2.7 文献综述（4）形状与动效",
    [
        {"text": "形状/图标研究", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  Bounding Box（矩形包围框）：Phan 2016、Wu 2024、Kim 2018 对照组", "size": 14},
        {"text": "·  Contact-Analog（与目标共形）：Chen 2024、Kim 2018", "size": 14},
        {"text": "·  Virtual Shadow / Dome+Tether：Kim 2016 / 2018（生态界面 EID）", "size": 14},
        {"text": "·  Stop Sign / Caution Sign：Kazazi 2015", "size": 14},
        {"text": "·  多元图标组合：Ma 2024 EID 'carpet'（梯形布局多组件）", "size": 14},
        {"text": " ", "size": 6},
        {"text": "动效设计研究", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  静态 vs 动态共形：多数研究偏好动态（Kim 2018、Wu 2024）", "size": 14},
        {"text": "·  闪烁注意捕获：Huo & Alla 2025（flashing 显著缩短反应时）", "size": 14},
        {"text": "·  跟随距离缩放：Kim 2018 virtual shadow 视角 2°–8° 动态", "size": 14},
        {"text": "·  Ma 2024 EID：颜色随风险增加从绿→黄→红渐变 + saturation 调整", "size": 14},
        {"text": " ", "size": 6},
        {"text": "→ 共识：动态共形优于静态屏幕固定；闪烁可补偿低紧迫感", "size": 16, "bold": True, "color": COLOR_PRIMARY},
    ],
    page=15, total=TOTAL
)

# Slide 16: Lit - Spatial Layout
make_richtext_slide(
    "2.8 文献综述（5）空间分布与 FOV",
    [
        {"text": "平面定位研究", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  Ye & Yin (2025, MDPI Electronics)：垂直面 vs 水平面 vs 混合面", "size": 14, "bold": True, "color": COLOR_ACCENT},
        {"text": "    → 水平面（路面投射）显著降低 inattentional blindness", "size": 14},
        {"text": "·  Kim 2018：单目 vs 体三维 HUD 对比 → 单目深度线索已足够", "size": 14},
        {"text": " ", "size": 6},
        {"text": "锁定方式研究", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  Wu et al. (2024)：BD（驾驶员视线）vs BR（路面投射）vs BW（行人跟随）", "size": 14},
        {"text": "    → BW（行人锁定）右转场景显著优于 BD/BR", "size": 14, "bold": True, "color": COLOR_ACCENT},
        {"text": "    → 首次注视时间：BW=616ms vs BD=2562ms vs BR=2729ms（p<0.001）", "size": 14},
        {"text": " ", "size": 6},
        {"text": "FOV 分级研究", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  Ma (2021)：65°（低速）/ 40°（高速）", "size": 14},
        {"text": "·  Teng (2023)：85° / 65° / 40°（速度三段）", "size": 14},
        {"text": "·  Zhong (2022)：FOV 12°×5° 用于颜色测量", "size": 14},
        {"text": " ", "size": 6},
        {"text": "→ 问题提出 3：空间锁定方式与 FOV 如何匹配不同驾驶场景？", "size": 16, "bold": True, "color": COLOR_ACCENT},
    ],
    page=16, total=TOTAL
)

# Slide 17: Lit - Individual Differences
make_richtext_slide(
    "2.9 文献综述（6）驾驶员个体差异",
    [
        {"text": "新手 vs 熟练驾驶员", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  Chen et al. (2024)：N=48 中国新手，contact-analog 在行人场景下反应更快", "size": 14},
        {"text": "·  Chen et al. (2025)：N=45 新手，多目标优先级显著优化反应时", "size": 14, "bold": True, "color": COLOR_ACCENT},
        {"text": "    （Hierarchical 模式：RT ↓、saccade counts ↓、gaze entropy ↓）", "size": 14},
        {"text": "·  Huo & Alla (2025)：新手对 AR 预警依赖度更高", "size": 14},
        {"text": " ", "size": 6},
        {"text": "老年 vs 年轻", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  Kazazi (2015)：老年组对 flow-point 触发感受不同；触发时机需个性化", "size": 14},
        {"text": " ", "size": 6},
        {"text": "环境差异", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  雾天/夜间：Zhang (2024) 中文版报告雾天 HUD 优势更显著", "size": 14},
        {"text": "·  城市/高速：FOV、TTC 阈值的最优值不同", "size": 14},
        {"text": " ", "size": 6},
        {"text": "→ 驾驶员经验是关键调节变量；本研究将聚焦『新手驾驶员』场景", "size": 16, "bold": True, "color": COLOR_PRIMARY},
    ],
    page=17, total=TOTAL
)

# Slide 18: Cognitive Load
make_richtext_slide(
    "2.10 文献综述（7）认知负荷与测量",
    [
        {"text": "行为指标", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  刹车反应时 SRT、最小 TTC、最大减速度、首次制动距离", "size": 14},
        {"text": " ", "size": 6},
        {"text": "眼动指标", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  首次注视时间 TTFF：警示捕获注意的速度（Wu 2024、Chen 2025）", "size": 14},
        {"text": "·  注视次数 Fixation count、扫视次数 Saccade count", "size": 14},
        {"text": "·  瞳孔直径：认知负荷的生理标志（Ma 2024）", "size": 14},
        {"text": "·  注视熵 Gaze Entropy：注意力分散程度（Chen 2025）", "size": 14},
        {"text": " ", "size": 6},
        {"text": "生理指标", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  EEG：N2/P3 波幅反映冲突监测（Strle 2023）", "size": 14},
        {"text": "·  HRV、GSR：唤醒水平", "size": 14},
        {"text": " ", "size": 6},
        {"text": "主观量表", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  NASA-TLX：6 维度心理负荷量表", "size": 14},
        {"text": "·  DALI：专用驾驶负荷量表", "size": 14},
        {"text": "·  SUS：系统可用性量表", "size": 14},
    ],
    page=18, total=TOTAL
)

# Slide 19: Research Gaps
make_richtext_slide(
    "2.11 研究空白与共识",
    [
        {"text": "已有共识", "size": 20, "bold": True, "color": COLOR_PRIMARY},
        {"text": " ", "size": 4},
        {"text": "  ✓ 共形增强 > 屏幕固定（Wu 2024、Kim 2018）", "size": 16, "color": RGBColor(0x20, 0x7E, 0x3D)},
        {"text": "  ✓ TTC = 2.5s 是常用警告触发点（Kim 2018、Lübbe 2017、Huo & Alla 2025）", "size": 16, "color": RGBColor(0x20, 0x7E, 0x3D)},
        {"text": "  ✓ 红色用于紧急警示，黄色用于注意提示", "size": 16, "color": RGBColor(0x20, 0x7E, 0x3D)},
        {"text": "  ✓ AR 预警显著提升新手驾驶员表现", "size": 16, "color": RGBColor(0x20, 0x7E, 0x3D)},
        {"text": " ", "size": 8},
        {"text": "研究空白", "size": 20, "bold": True, "color": COLOR_ACCENT},
        {"text": " ", "size": 4},
        {"text": "  ✗ 时间×空间维度的交互效应未系统研究", "size": 16},
        {"text": "  ✗ 多目标场景下的优先级设计研究极少（仅 Chen 2025 一项）", "size": 16},
        {"text": "  ✗ 中国驾驶员（特别是新手）的本土数据稀缺", "size": 16},
        {"text": "  ✗ 量化的、可被标准引用的设计参数表缺失", "size": 16},
        {"text": "  ✗ 多数研究为单一实验范式，缺乏跨场景验证", "size": 16},
        {"text": " ", "size": 8},
        {"text": "→ 本研究将聚焦上述空白，建立中国语境下的整合设计规范", "size": 16, "bold": True, "color": COLOR_PRIMARY},
    ],
    page=19, total=TOTAL
)

# Slide 20: Theory base
make_richtext_slide(
    "2.12 研究的理论基础",
    [
        {"text": "1. Multiple Resource Theory  多资源理论 (Wickens, 2002)", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "    AR-HUD 占用视觉空间-中央资源，但避免下视占用空间-外周资源", "size": 14},
        {"text": "    → 解释为何 HUD 优于 HDD（仪表盘下视）", "size": 14, "color": COLOR_TEXT_SOFT},
        {"text": " ", "size": 6},
        {"text": "2. Inattentional Blindness  非注意盲视 (Mack & Rock, 1998)", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "    非共形显示导致驾驶员忽视真实行人（『错觉感』）", "size": 14},
        {"text": "    → Ye & Yin (2025) 验证水平面共形减少 inattentional blindness", "size": 14, "color": COLOR_TEXT_SOFT},
        {"text": " ", "size": 6},
        {"text": "3. Contact-Analog Display Theory  接触类比显示 (Tönnis et al., 2007)", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "    与真实对象一一对应的虚拟显示降低认知映射成本", "size": 14},
        {"text": "    → 解释 Wu (2024) BW（行人锁定）优于 BD/BR（屏幕/路面）", "size": 14, "color": COLOR_TEXT_SOFT},
        {"text": " ", "size": 6},
        {"text": "4. PIEV Reaction Time Model  感知-识别-决策-行动模型 (Hooper, 1936)", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "    警告时机需为感知-识别-决策-行动提供足够时间", "size": 14},
        {"text": "    → 解释 TTC ≥ 2.5s 阈值的认知合理性", "size": 14, "color": COLOR_TEXT_SOFT},
    ],
    page=20, total=TOTAL
)

# Slide 21: TOC3
make_toc_slide("3")

# Slide 22: Section 3 title
make_section_title_slide("3", "研究问题", "Research Questions", total=TOTAL, page=22)

# Slide 23: Research questions summary
make_richtext_slide(
    "3.1 研究问题汇总",
    [
        {"text": "基于文献综述识别的三个核心问题：", "size": 20, "bold": True, "color": COLOR_PRIMARY},
        {"text": " ", "size": 8},
        {"text": "问题 1：AR-HUD 警告时机如何影响驾驶员对横穿行人的响应？", "size": 18, "bold": True, "color": COLOR_ACCENT},
        {"text": "  ·  不同 TTC 阈值（1.5s / 2.5s / 5.0s）的对比", "size": 14},
        {"text": "  ·  是否随驾驶员经验调整？（新手 vs 熟练）", "size": 14},
        {"text": " ", "size": 6},
        {"text": "问题 2：AR-HUD 空间锁定方式如何影响行人感知？", "size": 18, "bold": True, "color": COLOR_ACCENT},
        {"text": "  ·  屏幕固定 vs 行人锁定 vs 路面锁定", "size": 14},
        {"text": "  ·  在不同 FOV 与转向场景下的差异", "size": 14},
        {"text": " ", "size": 6},
        {"text": "问题 3：时间×空间设计因素是否存在交互效应？", "size": 18, "bold": True, "color": COLOR_ACCENT},
        {"text": "  ·  早警告 + 屏幕固定 vs 晚警告 + 行人锁定", "size": 14},
        {"text": "  ·  多目标场景下的优先级编码效用（拓展 Chen 2025）", "size": 14},
    ],
    page=23, total=TOTAL
)

# Slide 24: Research framework
make_richtext_slide(
    "3.2 研究框架",
    [
        {"text": "驾驶员因素（新手 vs 熟练）", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "                              ↓", "size": 16},
        {"text": "    ┌─── 时间设计（IV） ───┐    ┌─── 空间设计（IV） ───┐", "size": 14, "font": "Consolas"},
        {"text": "    │  TTC 阈值（1.5/2.5/5.0s） │    │  锁定方式（BD/BR/BW）   │", "size": 14, "font": "Consolas"},
        {"text": "    │  分级方式（无/2级/3级）   │    │  颜色编码（单色/分级）   │", "size": 14, "font": "Consolas"},
        {"text": "    │  持续时长                │    │  FOV / 形状 / 动效      │", "size": 14, "font": "Consolas"},
        {"text": "    └───────────────────────────┘    └───────────────────────────┘", "size": 14, "font": "Consolas"},
        {"text": "                              ↓", "size": 16},
        {"text": "    驾驶绩效 + 注意分配（DV）", "size": 18, "bold": True, "color": COLOR_ACCENT},
        {"text": "    ·  行为：刹车反应时 SRT、最小 TTC、最大减速度", "size": 14},
        {"text": "    ·  眼动：TTFF、注视次数、扫视次数、注视熵", "size": 14},
        {"text": "    ·  主观：NASA-TLX、SUS", "size": 14},
        {"text": "                              ↓", "size": 16},
        {"text": "    实验 1：时间维度  /  实验 2：空间维度  /  实验 3：交互×多目标", "size": 16, "bold": True, "color": COLOR_PRIMARY},
    ],
    page=24, total=TOTAL
)

# Slide 25: TOC4
make_toc_slide("4")

# Slide 26: Section 4 title
make_section_title_slide("4", "研究方案", "Research Plan", total=TOTAL, page=26)

# Slide 27: Plan overview
make_richtext_slide(
    "4.1 研究方案总览",
    [
        {"text": "本研究包含 3 项研究、4 个实验：", "size": 20, "bold": True, "color": COLOR_PRIMARY},
        {"text": " ", "size": 8},
        {"text": "研究一：警告时机研究", "size": 18, "bold": True, "color": COLOR_ACCENT},
        {"text": "·  实验 1：TTC 阈值 × 驾驶经验对驾驶绩效的影响", "size": 14},
        {"text": " ", "size": 8},
        {"text": "研究二：空间锁定与共形设计研究", "size": 18, "bold": True, "color": COLOR_ACCENT},
        {"text": "·  实验 2a：屏幕固定 vs 共形（驾驶员经验调节）", "size": 14},
        {"text": "·  实验 2b：FOV 分级与车速场景的匹配", "size": 14},
        {"text": " ", "size": 8},
        {"text": "研究三：时间×空间交互效应", "size": 18, "bold": True, "color": COLOR_ACCENT},
        {"text": "·  实验 3：多目标场景的优先级颜色编码（拓展 Chen 2025）", "size": 14},
        {"text": " ", "size": 8},
        {"text": "→ 总样本量约 210 人；预计 8 个月完成数据采集", "size": 16, "bold": True, "color": COLOR_PRIMARY},
    ],
    page=27, total=TOTAL
)

# Slide 28: Experiment 1
make_richtext_slide(
    "4.2 实验 1 — 警告时机研究",
    [
        {"text": "实验目的", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "  探究不同 TTC 阈值下 AR-HUD 警告对驾驶员制动反应的影响", "size": 14},
        {"text": " ", "size": 4},
        {"text": "研究方法", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "  ·  实验设计：3（TTC：1.5s / 2.5s / 4.0s）× 2（经验：新手 / 熟练）混合", "size": 14},
        {"text": "  ·  被试：60 名（g*power 计算，α=.05, 1-β=.80, f=.25）", "size": 14},
        {"text": "  ·  控制变量：车速 50km/h、行人触发距离、AR 图形（红色 bounding box）", "size": 14},
        {"text": "  ·  DV：刹车反应时 SRT、最小 TTC、最大减速度、首次制动距离", "size": 14},
        {"text": " ", "size": 4},
        {"text": "实验材料", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "  STISIM 驾驶模拟器 + AR-HUD 仿真（垂直虚像面，FOV 12°×5°）", "size": 14},
        {"text": "  10 试次 / 条件；单目标横穿行人场景", "size": 14},
        {"text": " ", "size": 4},
        {"text": "数据分析与预期", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "  2×3 混合 ANOVA + post-hoc", "size": 14},
        {"text": "  预期：TTC=2.5s 显著优于 1.5s（更早响应）和 4.0s（避免误警）；新手×短 TTC 交互显著", "size": 14, "color": COLOR_ACCENT},
    ],
    page=28, total=TOTAL
)

# Slide 29: Experiment 2a
make_richtext_slide(
    "4.3 实验 2a — 空间锁定方式与经验调节",
    [
        {"text": "实验目的", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "  探究不同锁定方式（屏幕固定 / 行人锁定 / 路面锁定）的认知负荷与绩效", "size": 14},
        {"text": " ", "size": 4},
        {"text": "研究方法（参照 Wu et al., 2024 设计）", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "  ·  实验设计：3（锁定：BD/BR/BW）× 2（经验：新手/熟练）混合", "size": 14},
        {"text": "  ·  被试：54 名", "size": 14},
        {"text": "  ·  DV：", "size": 14},
        {"text": "      行为：刹车反应时、最小 TTC", "size": 14},
        {"text": "      眼动：TTFF、注视次数、注视熵 Gaze Entropy", "size": 14},
        {"text": "      主观：NASA-TLX、SUS", "size": 14},
        {"text": " ", "size": 4},
        {"text": "实验材料", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "  HTC VIVE Pro Eye + Tobii Pro Lab 眼动 + Unity AR-HUD 仿真", "size": 14},
        {"text": "  场景：直行 / 左转 / 右转（参照 Wu 2024）", "size": 14},
        {"text": "  图形规格：红色 RGB(255,0,0) bounding box，TTC<3s 触发", "size": 14},
        {"text": " ", "size": 4},
        {"text": "预期：BW（行人锁定）在转向场景下显著优于 BD/BR；新手依赖度更高", "size": 14, "color": COLOR_ACCENT},
    ],
    page=29, total=TOTAL
)

# Slide 30: Experiment 2b
make_richtext_slide(
    "4.4 实验 2b — FOV 与车速场景匹配",
    [
        {"text": "实验目的", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "  探究最优 FOV 与车速的匹配关系", "size": 14},
        {"text": " ", "size": 4},
        {"text": "研究方法（参照 Ma 2021、Teng 2023）", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "  ·  实验设计：3（FOV：12° / 25° / 40°）× 3（车速：30 / 60 / 90 km/h）组内", "size": 14},
        {"text": "  ·  被试：36 名（重复测量设计）", "size": 14},
        {"text": "  ·  DV：识别正确率、首次注视时间、主观可用性", "size": 14},
        {"text": "  ·  控制：颜色、形状、锁定方式（均固定为 BW 行人锁定）", "size": 14},
        {"text": " ", "size": 4},
        {"text": "实验材料", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "  可调 FOV 的 AR-HUD 仿真平台（基于实验 2a 平台扩展）", "size": 14},
        {"text": " ", "size": 4},
        {"text": "数据分析与预期", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "  3×3 重复测量 ANOVA", "size": 14},
        {"text": "  预期：低速下窄 FOV 足够；高速需宽 FOV；FOV×车速交互显著", "size": 14, "color": COLOR_ACCENT},
        {"text": "  预期落点：低速 12° / 中速 25° / 高速 40°", "size": 14, "color": COLOR_ACCENT},
    ],
    page=30, total=TOTAL
)

# Slide 31: Experiment 3
make_richtext_slide(
    "4.5 实验 3 — 多目标优先级编码",
    [
        {"text": "实验目的（拓展 Chen et al., 2025）", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "  探究分级颜色优先级在多目标 AR-HUD 警告中的效用", "size": 14},
        {"text": "  扩展至中国驾驶员样本 + 增加目标数量维度", "size": 14},
        {"text": " ", "size": 4},
        {"text": "研究方法", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "  ·  实验设计：3（警告模式：Equivalent / Hierarchical / Baseline）× 2（目标数：2 / 3）混合", "size": 14},
        {"text": "  ·  被试：60 名新手驾驶员（驾龄 < 1 年）", "size": 14},
        {"text": "  ·  DV：反应时、saccade counts、gaze entropy、TTFF", "size": 14},
        {"text": " ", "size": 4},
        {"text": "实验材料", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "  AR 增强驾驶视频（参照 Chen 2025 范式）", "size": 14},
        {"text": "  Hierarchical：红=最高优先级，黄=次优先级，绿=低优先级", "size": 14},
        {"text": "  Equivalent：所有目标同色（红）", "size": 14},
        {"text": "  Baseline：无 AR 警告", "size": 14},
        {"text": " ", "size": 4},
        {"text": "预期：Hierarchical 在 3 目标场景下显著优于 Equivalent；复制 Chen 2025 结论", "size": 14, "color": COLOR_ACCENT},
    ],
    page=31, total=TOTAL
)

# Slide 32: Technical route
make_richtext_slide(
    "4.6 实验技术路线",
    [
        {"text": "阶段 1 — 准备期（M1-M3）", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "  ·  文献综述完善（已基本完成 - 40 篇核心文献）", "size": 14},
        {"text": "  ·  实验平台搭建（STISIM + HTC VIVE + Tobii）", "size": 14},
        {"text": "  ·  AR 警示刺激材料制作", "size": 14},
        {"text": "  ·  预实验（N=10）& 程序调试", "size": 14},
        {"text": " ", "size": 6},
        {"text": "阶段 2 — 数据收集（M4-M9）", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "  ·  实验 1（M4-M5）：60 被试", "size": 14},
        {"text": "  ·  实验 2a（M6）：54 被试", "size": 14},
        {"text": "  ·  实验 2b（M7）：36 被试", "size": 14},
        {"text": "  ·  实验 3（M8-M9）：60 被试", "size": 14},
        {"text": " ", "size": 6},
        {"text": "阶段 3 — 分析与写作（M10-M12）", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": "  ·  数据清洗与统计分析（SPSS / R / JASP）", "size": 14},
        {"text": "  ·  论文撰写与修改", "size": 14},
        {"text": "  ·  答辩准备", "size": 14},
    ],
    page=32, total=TOTAL
)

# Slide 33: Plan summary table
make_richtext_slide(
    "4.7 实验设计汇总",
    [
        {"text": "实验汇总表", "size": 18, "bold": True, "color": COLOR_PRIMARY},
        {"text": " ", "size": 6},
        {"text": "┌─────┬───────────┬─────────┬────┬───────────────────────┬──────────────────┐", "size": 12, "font": "Consolas"},
        {"text": "│ 实验 │   设计     │  IV    │ N  │     主要 DV           │   统计方法       │", "size": 12, "font": "Consolas", "bold": True},
        {"text": "├─────┼───────────┼─────────┼────┼───────────────────────┼──────────────────┤", "size": 12, "font": "Consolas"},
        {"text": "│  1  │ 3×2 混合  │ TTC×经验 │ 60 │ SRT, TTC, deceleration│ 2×3 混合 ANOVA   │", "size": 12, "font": "Consolas"},
        {"text": "│ 2a  │ 3×2 混合  │ 锁定×经验│ 54 │ SRT, TTFF, gaze ent.  │ 2×3 混合 ANOVA   │", "size": 12, "font": "Consolas"},
        {"text": "│ 2b  │ 3×3 组内  │ FOV×车速 │ 36 │ 识别率, TTFF          │ 3×3 RM-ANOVA     │", "size": 12, "font": "Consolas"},
        {"text": "│  3  │ 3×2 混合  │ 模式×目标│ 60 │ RT, saccade, entropy  │ 2×3 混合 ANOVA   │", "size": 12, "font": "Consolas"},
        {"text": "└─────┴───────────┴─────────┴────┴───────────────────────┴──────────────────┘", "size": 12, "font": "Consolas"},
        {"text": " ", "size": 6},
        {"text": "总样本：210 人（新手 ~120 + 熟练 ~90）", "size": 16, "bold": True, "color": COLOR_ACCENT},
        {"text": "总试次：约 14,400 次（210 × 平均 68 试次/被试）", "size": 14, "color": COLOR_TEXT_SOFT},
    ],
    page=33, total=TOTAL
)

# Slide 34: TOC5
make_toc_slide("5")

# Slide 35: Section 5 title
make_section_title_slide("5", "研究工作计划与进度安排", "Work Plan", total=TOTAL, page=35)

# Slide 36: Gantt chart
make_richtext_slide(
    "5.1 工作计划与进度安排（Gantt 图）",
    [
        {"text": " ", "size": 8},
        {"text": "                       M1 M2 M3 M4 M5 M6 M7 M8 M9 M10 M11 M12", "size": 13, "font": "Consolas", "bold": True},
        {"text": "文献综述完善           ■■■■", "size": 14, "font": "Consolas"},
        {"text": "实验平台搭建           ■■■■", "size": 14, "font": "Consolas"},
        {"text": "预实验                       ■■", "size": 14, "font": "Consolas"},
        {"text": "实验 1 数据采集                  ■■■■", "size": 14, "font": "Consolas", "color": COLOR_ACCENT},
        {"text": "实验 2a 数据采集                       ■■", "size": 14, "font": "Consolas", "color": COLOR_ACCENT},
        {"text": "实验 2b 数据采集                           ■■", "size": 14, "font": "Consolas", "color": COLOR_ACCENT},
        {"text": "实验 3 数据采集                              ■■■■", "size": 14, "font": "Consolas", "color": COLOR_ACCENT},
        {"text": "数据分析                             ■■■■■■■■■", "size": 14, "font": "Consolas"},
        {"text": "论文撰写                                 ■■■■■■■■", "size": 14, "font": "Consolas"},
        {"text": "论文修改 + 答辩准备                            ■■■■", "size": 14, "font": "Consolas"},
        {"text": " ", "size": 8},
        {"text": "里程碑节点：", "size": 16, "bold": True, "color": COLOR_PRIMARY},
        {"text": "  M3 末：完成预实验 + 平台验收", "size": 14},
        {"text": "  M9 末：完成所有数据采集", "size": 14},
        {"text": "  M11 末：论文初稿完成", "size": 14},
    ],
    page=36, total=TOTAL
)

# Slide 37: Expected outputs
make_richtext_slide(
    "5.2 预期产出",
    [
        {"text": "学术成果", "size": 20, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  1 篇 SSCI 二区论文（IJHCI / Ergonomics / Applied Ergonomics）", "size": 16},
        {"text": "·  1 篇 中文核心论文（《心理学报》 / 《人类工效学》）", "size": 16},
        {"text": "·  1 篇 会议论文（CHI / AutomotiveUI / HFES）", "size": 16},
        {"text": " ", "size": 8},
        {"text": "应用产出", "size": 20, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  AR-HUD 行人预警时空设计规范（草案）", "size": 16},
        {"text": "·  可复用的 AR 驾驶模拟实验平台", "size": 16},
        {"text": "·  4 项实验完整数据集（开源）", "size": 16},
        {"text": " ", "size": 8},
        {"text": "学术贡献", "size": 20, "bold": True, "color": COLOR_PRIMARY},
        {"text": "·  首次系统整合时间×空间维度的 AR-HUD 行人预警设计", "size": 16},
        {"text": "·  补充中国驾驶员（特别是新手）的本土证据", "size": 16},
        {"text": "·  扩展 Chen et al. (2025) 的多目标优先级研究框架", "size": 16},
    ],
    page=37, total=TOTAL
)

# Slide 38: Thanks
slide = add_blank_slide()
add_filled_rect(slide, 0, 0, prs.slide_width, prs.slide_height, COLOR_PRIMARY)
# Accent
add_filled_rect(slide, 0, Inches(3.4), prs.slide_width, Inches(0.08), COLOR_HIGHLIGHT)
add_text_box(slide, 0, Inches(2.5), prs.slide_width, Inches(1.5),
             "Thanks !", font_size=96, bold=True,
             color=COLOR_HIGHLIGHT, align=PP_ALIGN.CENTER,
             vertical=MSO_ANCHOR.MIDDLE, font_name="Calibri")
add_text_box(slide, 0, Inches(4.0), prs.slide_width, Inches(0.8),
             "敬请老师们提出宝贵意见", font_size=28, 
             color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
add_text_box(slide, 0, Inches(5.5), prs.slide_width, Inches(0.5),
             "致谢：导师、实验室、家人、所有 40 篇文献的作者", font_size=16,
             color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

prs.save(OUTPUT)
print(f"Saved {len(prs.slides)} slides to {OUTPUT}")
